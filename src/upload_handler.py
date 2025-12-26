"""
File upload and OCR processing handler
"""
import os
import io
from typing import Optional, Union
import easyocr
from PIL import Image



import fitz
import streamlit as st

class UploadHandler:
    def __init__(self):
        # Initialize EasyOCR reader (cached for performance)
        @st.cache_resource
        def get_ocr_reader():
            return easyocr.Reader(['en'], gpu=False)
        
        self.ocr_reader = get_ocr_reader()
        
    def extract_text_from_pdf(self, file_content: bytes) -> str:
        """Extract text from PDF using PyMuPDF (better than PyPDF2)"""
        try:
            pdf_document = fitz.open(stream=io.BytesIO(file_content), filetype="pdf")
            text_content = []
            
            for page_num in range(len(pdf_document)):
                page = pdf_document.load_page(page_num)
                text_content.append(page.get_text())
            
            pdf_document.close()
            return "\n".join(text_content)
        except Exception as e:
            print(f"PDF extraction error: {e}")
            return ""
    
    def extract_text_from_txt(self, file_content: bytes) -> str:
        """Extract text from TXT file"""
        try:
            return file_content.decode('utf-8', errors='ignore')
        except Exception as e:
            print(f"TXT extraction error: {e}")
            return ""
    

    def extract_text_from_image(self, file_content: bytes) -> str:
        """Extract text from image using EasyOCR"""
        try:
            # Convert bytes to PIL Image
            image = Image.open(io.BytesIO(file_content))
            
            # Use EasyOCR to extract text
            results = self.ocr_reader.readtext(image)
            
            # Extract text from results
            extracted_text = []
            for (bbox, text, confidence) in results:
                if confidence > 0.5:  # Only include confident extractions
                    extracted_text.append(text)
            
            return " ".join(extracted_text)
        except Exception as e:
            print(f"OCR extraction error: {e}")
            return ""
    
    def extract_text_from_pptx(self, file_content: bytes) -> str:
        """Extract text from PowerPoint (.pptx) files"""
        try:
            from pptx import Presentation
            
            presentation = Presentation(io.BytesIO(file_content))
            all_text = []
            
            for slide_idx, slide in enumerate(presentation.slides, 1):
                slide_text = [f"\n=== SLIDE {slide_idx} ==="]
                
                # Extract from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Identify if it's a title vs content
                        try:
                            if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                                if shape.placeholder_format.type == 1:  # Title
                                    slide_text.append(f"TITLE: {shape.text.strip()}")
                                else:
                                    slide_text.append(shape.text.strip())
                            else:
                                slide_text.append(shape.text.strip())
                        except:
                            # If we can't determine the type, just add the text
                            slide_text.append(shape.text.strip())
                    
                    # Extract bullet points if present
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text.strip():
                                slide_text.append(f"• {paragraph.text.strip()}")
                
                if len(slide_text) > 1:  # More than just the slide header
                    all_text.extend(slide_text)
            
            return "\n".join(all_text)
        except Exception as e:
            print(f"PPTX extraction error: {e}")
            return ""
    
    def process_upload(self, file_data: bytes, filename: str, user_id: str) -> tuple[str, str]:
        """
        Process uploaded file and return extracted text and file type
        
        Returns:
            tuple: (extracted_text, file_type)
        """
        # Determine file type
        file_extension = filename.lower().split('.')[-1]
        
        extracted_text = ""
        file_type = ""
        

        try:
            if file_extension == 'pdf':
                extracted_text = self.extract_text_from_pdf(file_data)
                file_type = "pdf"
            elif file_extension == 'pptx':
                extracted_text = self.extract_text_from_pptx(file_data)
                file_type = "powerpoint"
            elif file_extension in ['txt', 'md']:
                extracted_text = self.extract_text_from_txt(file_data)
                file_type = "text"
            elif file_extension in ['png', 'jpg', 'jpeg', 'bmp', 'tiff']:
                extracted_text = self.extract_text_from_image(file_data)
                file_type = "image"
            else:
                raise ValueError(f"Unsupported file type: {file_extension}")
            
            # Clean up extracted text
            extracted_text = self.clean_text(extracted_text)
            
            return extracted_text, file_type
            
        except Exception as e:
            print(f"Error processing file {filename}: {e}")
            return "", "error"
    
    def clean_text(self, text: str) -> str:
        """Clean and normalize extracted text"""
        if not text:
            return ""
        
        # Remove extra whitespace
        text = " ".join(text.split())
        
        # Remove very short lines (likely OCR noise) - reduced threshold from 3 to 2 to preserve more content
        lines = [line.strip() for line in text.split('\n') if len(line.strip()) > 2]
        
        return '\n'.join(lines)
    
    def chunk_text(self, text: str, chunk_size: int = 500, overlap: int = 50) -> list[str]:
        """
        Split text into overlapping chunks for better knowledge extraction
        
        Args:
            text: Input text to chunk
            chunk_size: Maximum characters per chunk
            overlap: Characters to overlap between chunks
            
        Returns:
            List of text chunks
        """
        if not text:
            return []
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + chunk_size
            
            # Try to break at sentence boundary
            if end < len(text):
                # Look for sentence endings near the end
                for i in range(end, max(start + chunk_size - 100, start), -1):
                    if text[i] in '.!?':
                        end = i + 1
                        break
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            # Move start position with overlap
            start = end - overlap
            
            # Ensure we make progress
            if start >= end:
                start = end
        
        return chunks

# Global handler instance
upload_handler = UploadHandler()
